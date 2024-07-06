import os
from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return text

def extract_text_from_folder(folder_path):
    with open("output.txt", "a") as file:
        for filename in os.listdir(folder_path):
            if filename.endswith(".pptx"):
                file_path = os.path.join(folder_path, filename)
                print(f"Extracting text from {filename}")
                slide_texts = extract_text_from_pptx(file_path)
                for slide_text in slide_texts:
                    # save to output file
                    file.write(slide_text + "\n")

# get current directory
folder_path = os.getcwd()
extract_text_from_folder(folder_path)
